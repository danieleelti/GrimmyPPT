import streamlit as st
import google.generativeai as genai
from google.oauth2 import service_account
from googleapiclient.discovery import build
from google.cloud import storage
import vertexai
from vertexai.preview.vision_models import ImageGenerationModel
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import os
import time
import uuid
import io

# --- CONFIGURAZIONE ---
st.set_page_config(page_title="Slide Monster: PRO COPYWRITER", page_icon="🦍", layout="wide")

# ======================================================
# ⚙️ RECUPERO DATI DA SECRETS
# ======================================================
if "slides_config" in st.secrets:
    DEF_TEMPLATE_ID = st.secrets["slides_config"]["template_id"]
    DEF_FOLDER_ID = st.secrets["slides_config"]["folder_id"]
else:
    DEF_TEMPLATE_ID = ""
    DEF_FOLDER_ID = ""

GCP_PROJECT_ID = "gen-lang-client-0247086002"
GCS_BUCKET_NAME = "bucket_grimmy"
GCP_LOCATION = "us-central1"

# --- GESTIONE STATO ---
if "app_state" not in st.session_state: st.session_state.app_state = "UPLOAD"
if "draft_data" not in st.session_state: st.session_state.draft_data = {}
if "final_images" not in st.session_state: st.session_state.final_images = {}
if "original_images" not in st.session_state: st.session_state.original_images = {} 

# --- INIZIALIZZAZIONE ---
try:
    if "gcp_service_account" in st.secrets and "json_content" in st.secrets["gcp_service_account"]:
        service_account_info = json.loads(st.secrets["gcp_service_account"]["json_content"])
    else:
        service_account_info = json.loads(st.secrets["GCP_SERVICE_ACCOUNT"])
    
    creds = service_account.Credentials.from_service_account_info(
        service_account_info,
        scopes=[
            'https://www.googleapis.com/auth/cloud-platform',
            'https://www.googleapis.com/auth/drive',
            'https://www.googleapis.com/auth/presentations'
        ]
    )

    genai.configure(api_key=st.secrets["GOOGLE_API_KEY"]) 
    drive_service = build('drive', 'v3', credentials=creds)
    slides_service = build('slides', 'v1', credentials=creds)
    
    vertexai.init(project=GCP_PROJECT_ID, location=GCP_LOCATION, credentials=creds)
    storage_client = storage.Client(credentials=creds, project=GCP_PROJECT_ID)
    bucket = storage_client.bucket(GCS_BUCKET_NAME)

except Exception as e:
    st.error(f"⚠️ Errore Inizializzazione: {e}")
    st.stop()

# ==========================================
# SIDEBAR
# ==========================================
with st.sidebar:
    st.header("🦍 Slide Monster PRO")
    
    with st.expander("⚙️ Configurazione Drive", expanded=True):
        tmpl = st.text_input("ID Template PPT", value=DEF_TEMPLATE_ID)
        fold = st.text_input("ID Cartella Output", value=DEF_FOLDER_ID)
        make_english = st.checkbox("🇬🇧 Genera anche versione Inglese", value=True)

    st.divider()

    st.subheader("🧠 Motore AI")
    # MODELLO UNICO E DEFINITIVO
    target_model = "models/gemini-3-pro-preview" 
    st.success(f"Running on: {target_model}")

    st.divider()
    if st.button("🔄 Reset Totale", type="secondary", use_container_width=True):
        st.session_state.app_state = "UPLOAD"
        st.session_state.draft_data = {}
        st.session_state.final_images = {}
        st.session_state.original_images = {}
        st.rerun()

# --- FUNZIONI CORE ---

def get_images_recursive_by_weight(shapes):
    """Cerca immagini ricorsivamente. Vince il PESO (Bytes)."""
    images_found = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                weight = len(blob) 
                images_found.append((weight, blob))
            except: pass
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            images_found.extend(get_images_recursive_by_weight(shape.shapes))
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                try:
                    blob = shape.image.blob
                    weight = len(blob)
                    images_found.append((weight, blob))
                except: pass
    return images_found

def analyze_pptx_content(file_obj):
    """Estrae testo e immagini (Heavyweight logic)."""
    prs = Presentation(file_obj)
    full_text = []
    extracted_images = {} 

    for i, slide in enumerate(prs.slides):
        s_txt = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                s_txt.append(shape.text.strip())
        visible_text = " | ".join(s_txt)

        notes_text = ""
        if slide.has_notes_slide:
            try:
                if slide.notes_slide.notes_text_frame:
                    notes_content = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_content:
                        notes_text = f"\n[[ ISTRUZIONI DALLE NOTE: {notes_content} ]]"
            except: pass
        
        full_text.append(f"SLIDE {i+1} CONTENUTO: {visible_text} {notes_text}")

        candidates = []
        candidates.extend(get_images_recursive_by_weight(slide.shapes))
        if slide.slide_layout:
            candidates.extend(get_images_recursive_by_weight(slide.slide_layout.shapes))
        if slide.slide_layout and slide.slide_layout.slide_master:
            candidates.extend(get_images_recursive_by_weight(slide.slide_layout.slide_master.shapes))
        
        if candidates:
            candidates.sort(key=lambda x: x[0], reverse=True)
            extracted_images[i] = candidates[0][1]
    
    return "\n---\n".join(full_text), extracted_images

def brain_process(text, model_name="models/gemini-3-pro-preview"):
    """
    PROMPT AGGRESSIVO E DETTAGLIATO - USING GEMINI 3 PRO PREVIEW
    """
    prompt = f"""
    Sei un SENIOR COPYWRITER esperto in Team Building e vendita di eventi B2B.
    Il tuo compito è analizzare il materiale grezzo (Slide + Note) e riscriverlo per VENDERE il format.
    
    ⚠️ REGOLE ASSOLUTE DI STILE:
    1. **NO EMOJI**. Sei professionale.
    2. **LUNGHEZZA:** I testi descrittivi (Pagina 2 e 3) devono essere CORPOSI (almeno 130-150 parole l'uno). Non fare riassuntini. Scrivi testi ricchi, ben articolati in paragrafi.
    3. **FORMATTAZIONE:** Poiché non puoi usare il grassetto, usa il **MAIUSCOLO** per enfatizzare le parole chiave e usa gli elenchi puntati (simbolo •) per dare ordine.
    4. **TONO:** Persuasivo, incoraggiante, emozionale ma concreto.
    
    ⚠️ ISTRUZIONI PER PAGINA:
    - **Cover:** Titolo del Format (Esatto) e Sottotitolo (Slogan).
    - **Pag 2 (L'Esperienza):** Descrivi la dinamica del gioco. Cosa succede? Come si interagisce? Scrivi TANTO testo. Fai immaginare al cliente di essere lì.
    - **Pag 3 (Il Valore):** Descrivi l'atmosfera, il coinvolgimento, l'energia. Perché questo format è unico? Scrivi TANTO testo.
    - **Pag 4 (Tecnica):** Svolgimento, Logistica, Tecnica. Qui devi essere un CHIRURGO. Elenchi puntati completi. Inserisci OGNI dettaglio tecnico trovato nel testo originale. Non tralasciare nulla.
    - **Pag 7 (Costi):** Crea una lista chiara e pulita di cosa è INCLUSO e cosa è ESCLUSO. Usa elenchi puntati.
    
    ⚠️ REGOLE LINGUA:
    1. Testi in **ITALIANO**.
    2. Prompt Immagini in **INGLESE**.
    
    STRUTTURA JSON:
    {{
        "page_1_cover": {{ 
            "title": "NOME DEL FORMAT", 
            "subtitle": "Slogan",
            "image_prompt": "Visual description in English" 
        }},
        "page_2_desc": {{ 
            "body": "Testo ESTESO (min 130 parole) sull'azione. Usa paragrafi e MAIUSCOLO per enfasi.", 
            "image_prompt": "Visual description in English" 
        }},
        "page_3_desc": {{ 
            "body": "Testo ESTESO (min 130 parole) sull'emozione. Usa paragrafi e MAIUSCOLO per enfasi.", 
            "image_prompt": "Visual description in English" 
        }},
        "page_4_details": {{
            "svolgimento": "Elenco puntato (•) DETTAGLIATO delle fasi.",
            "logistica": "Elenco puntato (•) DETTAGLIATO (spazi, tempi, pax).",
            "tecnica": "Elenco puntato (•) DETTAGLIATO (audio, video, prese)."
        }},
        "page_7_costi": {{
            "dettaglio": "Elenco puntato (•) CHIARO: IL COSTO INCLUDE... / IL COSTO NON COMPRENDE..."
        }}
    }}
    """
    model = genai.GenerativeModel(model_name)
    try:
        resp = model.generate_content(f"{prompt}\n\nTESTO SORGENTE:\n{text}", generation_config={"response_mime_type": "application/json"})
        return json.loads(resp.text)
    except Exception as e:
        st.error(f"Errore Gemini Brain: {e}")
        return None

def translate_struct_to_english(ai_data):
    """Traduce la struttura mantenendo la formattazione - USING GEMINI 3 PRO PREVIEW"""
    prompt = """
    You are a professional translator and copywriter. 
    Translate the values in the following JSON from Italian to English.
    
    RULES:
    1. **KEEP THE FORMAT NAME (TITLE) AS IS**. Do not translate proper names.
    2. **Translate fully**. Do not summarize. Keep the text long and persuasive.
    3. **Maintain formatting**: Keep the bullet points (•) and UPPERCASE words for emphasis.
    4. Do not translate keys or 'image_prompt'.
    
    Return ONLY valid JSON.
    """
    # Forziamo Gemini 3 Pro anche qui
    model = genai.GenerativeModel("models/gemini-3-pro-preview") 
    try:
        resp = model.generate_content(f"{prompt}\n\nJSON:\n{json.dumps(ai_data)}", generation_config={"response_mime_type": "application/json"})
        return json.loads(resp.text)
    except Exception as e:
        print(f"Errore Traduzione Struct: {e}")
        return ai_data 

def get_template_static_text(presentation_id):
    try:
        prs = slides_service.presentations().get(presentationId=presentation_id).execute()
        texts_to_translate = set()
        for slide in prs.get('slides', []):
            for el in slide.get('pageElements', []):
                if 'shape' in el and 'text' in el['shape']:
                    for tr in el['shape']['text']['textElements']:
                        if 'textRun' in tr and 'content' in tr['textRun']:
                            content = tr['textRun']['content'].strip()
                            if content and "{{" not in content and "}}" not in content and len(content) > 2:
                                texts_to_translate.add(content)
        return list(texts_to_translate)
    except: return []

def translate_list_strings(text_list):
    if not text_list: return {}
    prompt = "Translate these Italian strings to English for a corporate presentation. Return JSON {original: translation}."
    # Forziamo Gemini 3 Pro anche qui
    model = genai.GenerativeModel("models/gemini-3-pro-preview")
    try:
        resp = model.generate_content(f"{prompt}\n\nLIST:\n{json.dumps(text_list)}", generation_config={"response_mime_type": "application/json"})
        return json.loads(resp.text)
    except: return {}

def apply_static_translations(presentation_id, translation_map):
    if not translation_map: return
    reqs = []
    for it, en in translation_map.items():
        if it and en and it != en:
            reqs.append({'replaceAllText': {'containsText': {'text': it, 'matchCase': True}, 'replaceText': en}})
    if reqs:
        chunk_size = 50
        for i in range(0, len(reqs), chunk_size):
            try: slides_service.presentations().batchUpdate(presentationId=presentation_id, body={'requests': reqs[i:i+chunk_size]}).execute()
            except: pass

def upload_bytes_to_bucket(image_bytes):
    try:
        filename = f"img_{uuid.uuid4()}.png"
        blob = bucket.blob(filename)
        blob.upload_from_string(image_bytes, content_type="image/png")
        blob.make_public() # FONDAMENTALE PER VISIBILITÀ
        return f"https://storage.googleapis.com/{GCS_BUCKET_NAME}/{filename}"
    except Exception as e:
        st.error(f"Errore Upload Bucket: {e}")
        return None

def generate_imagen_safe(prompt):
    for i in range(3):
        try:
            model = ImageGenerationModel.from_pretrained("imagen-3.0-generate-001")
            images = model.generate_images(prompt=prompt, number_of_images=1, aspect_ratio="16:9", person_generation="allow_adult")
            if images: return images[0]._image_bytes
        except Exception as e:
            if "429" in str(e) or "Quota" in str(e): time.sleep(10 * (i+1))
            else: return None
    return None

def find_image_element_id_smart(prs_id, label):
    try:
        prs = slides_service.presentations().get(presentationId=prs_id).execute()
        label_clean = label.strip().upper()
        for slide in prs.get('slides', []):
            for el in slide.get('pageElements', []):
                if 'description' in el and el['description'].strip().upper() == label_clean:
                    return el['objectId']
    except: pass
    return None

def worker_bot_finalize(template_id, folder_id, filename, ai_data, urls_map, translate_mode=False):
    try:
        copy = drive_service.files().copy(
            fileId=template_id, body={'name': filename, 'parents': [folder_id]}, supportsAllDrives=True
        ).execute()
        new_id = copy.get('id')
        
        final_data = ai_data
        
        # LOGICA TRADUZIONE
        if translate_mode:
            st.toast(f"🇬🇧 Traduzione AI in corso (Gemini 3 Pro): {filename}")
            # Traduciamo i dati dinamici
            final_data = translate_struct_to_english(ai_data)
            
            # Controllo sicurezza: se la traduzione è identica all'originale, qualcosa non va
            if final_data['page_2_desc']['body'] == ai_data['page_2_desc']['body']:
                 st.warning(f"⚠️ Attenzione: La traduzione inglese di {filename} sembra identica all'italiano. Riprovo...")
                 final_data = translate_struct_to_english(ai_data) # Retry once
            
            # Traduciamo i testi statici del template
            static_texts = get_template_static_text(new_id)
            if static_texts:
                t_map = translate_list_strings(static_texts)
                apply_static_translations(new_id, t_map)
        
        # --- TITOLO GLOBALE (PUNTO 1 RISOLTO) ---
        main_format_title = final_data.get('page_1_cover', {}).get('title', 'Format')

        reqs = []
        # Sostituzione globale di {{TITLE}} su tutte le slide
        reqs.append({'replaceAllText': {'containsText': {'text': '{{TITLE}}'}, 'replaceText': main_format_title}})
        
        # Cover Subtitle
        if 'page_1_cover' in final_data:
            reqs.append({'replaceAllText': {'containsText': {'text': '{{SUBTITLE}}'}, 'replaceText': final_data['page_1_cover'].get('subtitle', '')}})
        
        # Desc 1
        if 'page_2_desc' in final_data:
            reqs.append({'replaceAllText': {'containsText': {'text': '{{BODY_1}}'}, 'replaceText': final_data['page_2_desc'].get('body', '')}})
        
        # Desc 2
        if 'page_3_desc' in final_data:
            reqs.append({'replaceAllText': {'containsText': {'text': '{{BODY_2}}'}, 'replaceText': final_data['page_3_desc'].get('body', '')}})
        
        # Dettagli Tecnici
        if 'page_4_details' in final_data:
            reqs.append({'replaceAllText': {'containsText': {'text': '{{SVOLGIMENTO}}'}, 'replaceText': final_data['page_4_details'].get('svolgimento', '')}})
            reqs.append({'replaceAllText': {'containsText': {'text': '{{LOGISTICA}}'}, 'replaceText': final_data['page_4_details'].get('logistica', '')}})
            reqs.append({'replaceAllText': {'containsText': {'text': '{{TECNICA}}'}, 'replaceText': final_data['page_4_details'].get('tecnica', '')}})
        
        # Costi
        if 'page_7_costi' in final_data:
            reqs.append({'replaceAllText': {'containsText': {'text': '{{DETTAGLIO_COSTO}}'}, 'replaceText': final_data['page_7_costi'].get('dettaglio', '')}})

        if reqs:
            slides_service.presentations().batchUpdate(presentationId=new_id, body={'requests': reqs}).execute()

        for label, url in urls_map.items():
            if url:
                el_id = find_image_element_id_smart(new_id, label)
                if el_id:
                    req = {'replaceImage': {'imageObjectId': el_id, 'imageReplaceMethod': 'CENTER_CROP', 'url': url}}
                    try: slides_service.presentations().batchUpdate(presentationId=new_id, body={'requests': [req]}).execute()
                    except: pass
        return new_id
    except Exception as e:
        print(f"Error in finalize: {e}")
        return None

# ==========================================
# MAIN INTERFACE
# ==========================================
st.title("🦍 Slide Monster: PRO COPYWRITER")

# --- FASE 1: UPLOAD ---
if st.session_state.app_state == "UPLOAD":
    st.markdown("### 1. Carica le presentazioni")
    with st.container(border=True):
        uploaded = st.file_uploader("PPTX", accept_multiple_files=True, type=['pptx'])
        col_act1, col_act2 = st.columns([1, 4])
        with col_act1:
            if st.button("🧠 ANALIZZA", type="primary", use_container_width=True):
                if uploaded:
                    st.session_state.draft_data = {}
                    st.session_state.final_images = {}
                    st.session_state.original_images = {}
                    
                    bar = st.progress(0)
                    for i, f in enumerate(uploaded):
                        fname = f.name.replace(".pptx", "") + "_ITA"
                        txt, imgs_dict = analyze_pptx_content(f)
                        data = brain_process(txt) # Default a Gemini 3 Pro
                        
                        if data:
                            st.session_state.draft_data[fname] = {"ai_data": data}
                            st.session_state.final_images[fname] = {}
                            st.session_state.original_images[fname] = imgs_dict
                        bar.progress((i+1)/len(uploaded))
                    st.session_state.app_state = "EDIT"
                    st.rerun()
        with col_act2:
            st.caption("Analisi Pro: Testi lunghi, formattazione elenchi, no emoji, traduzione Inglese potenziata.")

# --- FASE 2: EDITING ---
elif st.session_state.app_state == "EDIT":
    
    col_h1, col_h2 = st.columns([3, 1])
    with col_h1:
        st.info("✏️ **Sala di Regia**: Controlla i contenuti. I testi sono stati generati per essere 'Corporate' e 'Vendevoli'.")
    with col_h2:
        if st.button("💾 SALVA SU DRIVE", type="primary", use_container_width=True):
            bar = st.progress(0)
            total_ops = len(st.session_state.draft_data)
            
            for i, (fname, content) in enumerate(st.session_state.draft_data.items()):
                url_map = {}
                saved = st.session_state.final_images.get(fname, {})
                
                if 'cover' in saved: url_map['IMG_1'] = saved['cover'] 
                if 'desc_1' in saved: url_map['IMG_2'] = saved['desc_1'] 
                if 'desc_2' in saved: url_map['IMG_3'] = saved['desc_2'] 
                
                # ITA
                st.toast(f"🇮🇹 Saving ITA: {fname}")
                res_ita = worker_bot_finalize(tmpl, fold, fname, content['ai_data'], url_map, translate_mode=False)
                
                # ENG
                if make_english:
                    if "_ITA" in fname:
                        fname_eng = fname.replace("_ITA", "_ENG")
                    else:
                        fname_eng = fname + "_ENG"
                        
                    st.toast(f"🇬🇧 Saving ENG: {fname_eng}")
                    res_eng = worker_bot_finalize(tmpl, fold, fname_eng, content['ai_data'], url_map, translate_mode=True)

                if res_ita: st.success(f"✅ Fatto: {fname}")
                else: st.error(f"❌ Errore: {fname}")
                
                bar.progress((i+1)/total_ops)
            st.balloons()
            time.sleep(2)

    for fname, content in st.session_state.draft_data.items():
        data = content['ai_data']
        orig_imgs = st.session_state.original_images.get(fname, {})
        
        st.markdown(f"## 📂 {fname}")
        
        # AGGIUNTO IL TAB COSTI
        tabs = st.tabs(["🏠 1. Cover", "📄 2. L'Esperienza", "📄 3. L'Emozione", "🛠️ 4. Scheda Tecnica", "💰 7. Costi"])
        
        # --- TAB 1: COVER ---
        with tabs[0]:
            c1, c2, c3 = st.columns([1.5, 1, 1])
            with c1:
                st.markdown("#### Testo Cover")
                st.caption("Questo titolo andrà su TUTTE le slide (sostituisce {{TITLE}})")
                new_t = st.text_input("Titolo Format (Esatto)", value=data['page_1_cover'].get('title', ''), key=f"t1_{fname}")
                new_s = st.text_input("Sottotitolo", value=data['page_1_cover'].get('subtitle', ''), key=f"s1_{fname}")
                st.session_state.draft_data[fname]['ai_data']['page_1_cover']['title'] = new_t
                st.session_state.draft_data[fname]['ai_data']['page_1_cover']['subtitle'] = new_s
            with c2:
                st.markdown("#### AI Image")
                p = st.text_area("Prompt", value=data['page_1_cover'].get('image_prompt', ''), height=70, key=f"p1_{fname}")
                if st.button("Genera", key=f"b1_{fname}"):
                    bytes_img = generate_imagen_safe(p)
                    if bytes_img: 
                        st.session_state.final_images[fname]['cover'] = upload_bytes_to_bucket(bytes_img)
                        st.rerun() # Refresh immediato
                
                if st.session_state.final_images[fname].get('cover'):
                    st.image(st.session_state.final_images[fname]['cover'], width=200)
                    st.success("Immagine Pronta")
            with c3:
                st.markdown("#### Originale")
                if orig_imgs.get(0):
                    st.image(orig_imgs[0], width=200, caption=f"Peso: {len(orig_imgs[0])//1024} KB")
                    if st.button("Usa Originale", key=f"bo1_{fname}"):
                        st.session_state.final_images[fname]['cover'] = upload_bytes_to_bucket(orig_imgs[0]); st.rerun()

        # --- TAB 2: DESC 1 ---
        with tabs[1]:
            c1, c2, c3 = st.columns([1.5, 1, 1])
            with c1:
                st.markdown("#### L'Esperienza")
                st.info(f"Titolo Slide: {data['page_1_cover'].get('title', '')}")
                new_b = st.text_area("Body 1 (Lungo e Formattato)", value=data['page_2_desc'].get('body', ''), height=300, key=f"b2_{fname}")
                st.session_state.draft_data[fname]['ai_data']['page_2_desc']['body'] = new_b
            with c2:
                st.markdown("#### AI Image")
                p = st.text_area("Prompt", value=data['page_2_desc'].get('image_prompt', ''), height=70, key=f"p2_{fname}")
                if st.button("Genera", key=f"b2_gen_{fname}"):
                    bytes_img = generate_imagen_safe(p)
                    if bytes_img: 
                        st.session_state.final_images[fname]['desc_1'] = upload_bytes_to_bucket(bytes_img)
                        st.rerun()
                if st.session_state.final_images[fname].get('desc_1'):
                    st.image(st.session_state.final_images[fname]['desc_1'], width=200)
                    st.success("Immagine Pronta")
            with c3:
                st.markdown("#### Originale")
                if orig_imgs.get(1):
                    st.image(orig_imgs[1], width=200, caption=f"Peso: {len(orig_imgs[1])//1024} KB")
                    if st.button("Usa Originale", key=f"bo2_{fname}"):
                        st.session_state.final_images[fname]['desc_1'] = upload_bytes_to_bucket(orig_imgs[1]); st.rerun()

        # --- TAB 3: DESC 2 ---
        with tabs[2]:
            c1, c2, c3 = st.columns([1.5, 1, 1])
            with c1:
                st.markdown("#### Il Valore")
                st.info(f"Titolo Slide: {data['page_1_cover'].get('title', '')}")
                new_b = st.text_area("Body 2 (Lungo e Formattato)", value=data['page_3_desc'].get('body', ''), height=300, key=f"b3_{fname}")
                st.session_state.draft_data[fname]['ai_data']['page_3_desc']['body'] = new_b
            with c2:
                st.markdown("#### AI Image")
                p = st.text_area("Prompt", value=data['page_3_desc'].get('image_prompt', ''), height=70, key=f"p3_{fname}")
                if st.button("Genera", key=f"b3_gen_{fname}"):
                    bytes_img = generate_imagen_safe(p)
                    if bytes_img: 
                        st.session_state.final_images[fname]['desc_2'] = upload_bytes_to_bucket(bytes_img)
                        st.rerun()
                if st.session_state.final_images[fname].get('desc_2'):
                    st.image(st.session_state.final_images[fname]['desc_2'], width=200)
                    st.success("Immagine Pronta")
            with c3:
                st.markdown("#### Originale")
                if orig_imgs.get(2):
                    st.image(orig_imgs[2], width=200, caption=f"Peso: {len(orig_imgs[2])//1024} KB")
                    if st.button("Usa Originale", key=f"bo3_{fname}"):
                        st.session_state.final_images[fname]['desc_2'] = upload_bytes_to_bucket(orig_imgs[2]); st.rerun()

        # --- TAB 4: DETAILS ---
        with tabs[3]:
            st.markdown("#### 🛠️ Dettagli Tecnici")
            c1, c2, c3 = st.columns(3)
            with c1:
                st.markdown("**Svolgimento**")
                v1 = st.text_area("Testo", value=data['page_4_details'].get('svolgimento', ''), height=400, key=f"d1_{fname}")
                st.session_state.draft_data[fname]['ai_data']['page_4_details']['svolgimento'] = v1
            with c2:
                st.markdown("**Logistica**")
                v2 = st.text_area("Testo", value=data['page_4_details'].get('logistica', ''), height=400, key=f"d2_{fname}")
                st.session_state.draft_data[fname]['ai_data']['page_4_details']['logistica'] = v2
            with c3:
                st.markdown("**Tecnica**")
                v3 = st.text_area("Testo", value=data['page_4_details'].get('tecnica', ''), height=400, key=f"d3_{fname}")
                st.session_state.draft_data[fname]['ai_data']['page_4_details']['tecnica'] = v3

        # --- TAB 5: COSTI ---
        with tabs[4]:
            st.markdown("#### 💰 Dettagli Economici (Slide 7)")
            det = st.text_area("Dettaglio Costi (Include/Esclude)", value=data.get('page_7_costi', {}).get('dettaglio', ''), height=400, key=f"c_det_{fname}")
            
            if 'page_7_costi' not in st.session_state.draft_data[fname]['ai_data']:
                st.session_state.draft_data[fname]['ai_data']['page_7_costi'] = {}
            st.session_state.draft_data[fname]['ai_data']['page_7_costi']['dettaglio'] = det

        st.markdown("---")
