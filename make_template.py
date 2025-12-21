from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_grimmy_template():
    prs = Presentation() # Crea un PPT vuoto standard (16:9 solitamente)
    
    # Pulizia: Rinominiamo i layout standard con i nomi che vuole Grimmy
    # I layout standard di solito sono 11. Ne usiamo alcuni.
    
    master = prs.slide_master
    layouts = master.slide_layouts
    
    # MAPPATURA: (Indice originale) -> (Nome per Grimmy)
    # 0 = Title Slide
    # 1 = Title and Content
    # 2 = Section Header
    # 3 = Two Content
    # 4 = Comparison
    # 5 = Title Only
    # 6 = Blank
    # ...
    
    # 1. COVER MAIN (Uso il layout 0)
    layout_cover = layouts[0]
    layout_cover.name = "Cover_Main"
    # Aggiungi nota per te
    add_hint(layout_cover, "COVER: Qui Grimmy metterà Titolo, Sottotitolo e l'Immagine NanoBanana")

    # 2. INTRO CONCEPT (Uso il layout 1)
    layout_intro = layouts[1]
    layout_intro.name = "Intro_Concept"
    add_hint(layout_intro, "INTRO: Qui andrà il concept emotivo")

    # 3. ACTIVITY DETAIL (Uso il layout 3 - Due Contenuti, ottimo per testo + foto)
    layout_activity = layouts[3]
    layout_activity.name = "Activity_Detail"
    add_hint(layout_activity, "ACTIVITY: Qui Grimmy mette dettagli operativi e foto")

    # 4. TECHNICAL GRID (Uso il layout 4 - Comparison)
    layout_tech = layouts[4]
    layout_tech.name = "Technical_Grid"
    add_hint(layout_tech, "TECHNICAL: Durata, pax, e scheda tecnica")

    # 5. LOGISTICS INFO (Uso il layout 2 - Section Header, modificato)
    layout_logistics = layouts[2]
    layout_logistics.name = "Logistics_Info"
    add_hint(layout_logistics, "LOGISTICS: Cosa è incluso/escluso")

    # --- PAGINE FISSE (STANDARD) ---
    # Usiamo altri layout generici e li rinominiamo.
    
    # 6. Standard Training (Layout 5)
    l_train = layouts[5]
    l_train.name = "Standard_Training"
    add_hint(l_train, "FISSO: Incolla qui la grafica 'Formazione' definitiva.")

    # 7. Standard Extras (Layout 6 - Blank)
    l_extra = layouts[6]
    l_extra.name = "Standard_Extras"
    add_hint(l_extra, "FISSO: Incolla qui grafica Gadget/Foto/Video.")

    # 8. Standard Payment (Layout 7 - Content with Caption)
    # Se non esiste l'indice 7, usiamo il 6 duplicato concettualmente (qui forziamo il nome su uno esistente)
    if len(layouts) > 7:
        l_pay = layouts[7]
        l_pay.name = "Standard_Payment"
        add_hint(l_pay, "FISSO: Dati bancari e IBAN.")
    else:
        # Fallback se il tema è strano
        layouts[6].name = "Standard_Payment" 

    # 9. Closing Contact (Layout 8)
    if len(layouts) > 8:
        l_close = layouts[8]
        l_close.name = "Closing_Contact"
        add_hint(l_close, "FISSO: Contatti finali.")

    # Salvataggio
    output_name = "Template_Grimmy_Base.pptx"
    prs.save(output_name)
    print(f"✅ Fatto! File creato: {output_name}")
    print("Ora aprilo in PowerPoint e applica la tua grafica (Colori, Font, Loghi) SENZA cambiare i nomi dei layout.")

def add_hint(layout, text):
    """Aggiunge una casella di testo nel master per ricordarti cosa va dove"""
    txBox = layout.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = text
    tf.paragraphs[0].font.color.rgb = None # Default
    tf.paragraphs[0].font.size = Pt(10)

if __name__ == "__main__":
    create_grimmy_template()
