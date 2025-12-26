import streamlit as st
import google.generativeai as genai
import json
import requests
import io
from pptx.util import Inches, Pt

# --- 1. ANALISI TESTO (Scenario Emozionale) ---
def analyze_content(context, gemini_model):
    try:
        model = genai.GenerativeModel(gemini_model)
        prompt_text = f"""
        Sei un Copywriter creativo. Stiamo scrivendo la PAGINA 2: LO SCENARIO (CONCEPT).
        
        Analizza il testo e restituisci:
        1. "format_name": Il nome del format (Titolo slide).
        2. "emotional_text": Un testo EVOCATIVO, EMOZIONALE e ISPIRAZIONALE che descriva l'atmosfera del format. (Sarà scritto in corsivo nel template, quindi usa un tono elegante e coinvolgente). Max 2-3 frasi.
        3. "imagen_prompt": Un prompt per un'immagine di SFONDO A TUTTA PAGINA. Deve essere atmosferica, ampia e cinematografica.

        RISPONDI SOLO JSON: 
        {{
            "format_name": "...", 
            "emotional_text": "...", 
            "imagen_prompt": "..."
        }}

        TESTO SORGENTE: {context[:6000]}
        """
        res_text = model.generate_content(prompt_text, generation_config={"response_mime_type": "application/json"})
        return json.loads(res_text.text)
    except Exception as e:
        st.error(f"Errore Analisi Page 2: {e}")
        return None

# --- 2. GENERAZIONE IMMAGINE ---
def generate_image(prompt, api_key, model_name):
    if not model_name.startswith("models/"): model_name = f"models/{model_name}"
    url = f"https://generativelanguage.googleapis.com/v1beta/{model_name}:predict?key={api_key}"
    headers = {"Content-Type": "application/json"}
    data = {"instances": [{"prompt": prompt}], "parameters": {"aspectRatio": "16:9", "sampleCount": 1}}
    
    try:
        response = requests.post(url, headers=headers, json=data)
        response.raise_for_status()
        result = response.json()
        if "predictions" in result:
            import base64
            return base64.b64decode(result["predictions"][0]["bytesBase64Encoded"])
        return None
    except Exception as e:
        st.error(f"Errore Imagen Page 2: {e}")
        return None

# --- 3. INSERIMENTO NELLA SLIDE (Logica "Search & Rescue") ---
def insert_into_slide(slide, data, img_bytes):
    try:
        # A. TITOLO (Cerca shape Titolo o usa la prima shape in alto con testo)
        if slide.shapes.title:
            slide.shapes.title.text = data.get("format_name", "")
        
        # B. TESTO EMOZIONALE (Caccia al Tesoro)
        target_shape = None
        candidates = []
        
        # 1. Cerca nei Placeholders ufficiali (escludendo il titolo)
        for shape in slide.placeholders:
            if shape.has_text_frame and shape != slide.shapes.title:
                candidates.append(shape)
        
        # 2. Se vuoto, cerca in TUTTE le Shapes (es. box di testo statici nel layout)
        if not candidates:
            # Cerchiamo shape che hanno testo e non sono il titolo
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    # Filtro anti-rumore: Ignoriamo cose minuscole (es. numeri pagina)
                    if shape.width > Inches(2) and shape.height > Inches(0.5):
                        candidates.append(shape)

        # 3. Selezione del Vincitore
        if candidates:
            # Prendiamo il box PIÙ GRANDE (per area), assumendo che sia il corpo del testo
            target_shape = max(candidates, key=lambda s: s.width * s.height)
            
            # Scrittura
            tf = target_shape.text_frame
            tf.clear() # Pulisce "Edit Text"
            p = tf.paragraphs[0]
            p.text = data.get("emotional_text", "")
            # Nota: Manteniamo la formattazione originale del template (corsivo, font, ecc.)
            st.toast(f"Testo inserito in un box esistente (Area: {int(target_shape.width*target_shape.height)})", icon="✅")
        
        else:
            # 4. ULTIMO FALLBACK: CREAZIONE MANUALE
            # Se non esiste nessun box, ne creiamo uno noi bello grande
            st.warning("⚠️ Nessun box di testo trovato nel template. Ne creo uno nuovo.")
            left = Inches(1)
            top = Inches(2)
            width = Inches(10)
            height = Inches(3)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            p = textbox.text_frame.add_paragraph()
            p.text = data.get("emotional_text", "")
            p.font.size = Pt(24) # Font leggibile
            p.font.italic = True # Forziamo il corsivo richiesto
            
            target_shape = textbox # Per riferimento Z-Order dopo

        # C. IMMAGINE DI SFONDO (Z-Order Fix)
        if img_bytes:
            image_stream = io.BytesIO(img_bytes)
            
            # Inserisci a tutto schermo
            pic = slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
            
            # Sposta DIETRO (Send to Back)
            try:
                slide.shapes._spTree.remove(pic._element)
                # Indice 2 è solitamente sicuro per stare dietro al testo ma non rompere il file
                slide.shapes._spTree.insert(2, pic._element)
            except Exception:
                pass

        return True
    except Exception as e:
        st.error(f"Errore scrittura Page 2: {e}")
        return False
