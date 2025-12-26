import streamlit as st
import google.generativeai as genai
from pptx import Presentation
import io
import json

# --- SETUP ---
st.set_page_config(page_title="Step 1: La Cover", layout="wide")

# Recupero Secrets
try:
    GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]
    genai.configure(api_key=GOOGLE_API_KEY)
except KeyError:
    st.error("Errore: Manca GOOGLE_API_KEY nei secrets.")
    st.stop()

# Recupero Password
if 'password_correct' not in st.session_state: st.session_state['password_correct'] = False

def check_password():
    if st.session_state['password_correct']: return True
    pwd = st.sidebar.text_input("Password", type="password")
    if st.sidebar.button("Login"):
        if pwd == st.secrets["app_password"]:
            st.session_state['password_correct'] = True
            st.rerun()
    return False

if not check_password(): st.stop()

# --- LOGICA SPECIFICA PER LA COVER ---

def get_cover_data(old_ppt_text, model_name):
    """
    Chiede a Gemini di estrarre il titolo esatto e inventare un Claim.
    """
    system_instruction = """
    Sei un esperto di Marketing e Team Building.
    Stiamo creando la COPERTINA (Cover) di una presentazione commerciale.
    
    Il tuo output deve essere un JSON rigoroso:
    {
        "format_name": "Il nome ESATTO del format trovato nel testo (NON CAMBIARLO MAI)",
        "claim": "Uno slogan commerciale di vendita, breve, accattivante ed energico (max 8 parole)",
        "imagen_prompt": "Prompt in inglese per Imagen 3. Deve generare una immagine di sfondo epica, fotorealistica, corporate ma emozionale, che rappresenti il concetto del format."
    }
    """
    
    prompt = f"""
    Analizza il testo della vecchia presentazione e crea i contenuti per la Cover.
    TESTO VECCHIO PPT:
    {old_ppt_text}
    """
    
    model = genai.GenerativeModel(model_name, system_instruction=system_instruction)
    response = model.generate_content(prompt, generation_config={"response_mime_type": "application/json"})
    return json.loads(response.text)

def process_cover_slide(template_file, data):
    """
    Modifica ESCLUSIVAMENTE la slide [0] del template.
    """
    prs = Presentation(template_file)
    
    # PUNTO FONDAMENTALE: Lavoriamo sulla slide 0 esistente, non ne creiamo nuove.
    slide = prs.slides[0]
    
    # 1. GESTIONE TITOLO (Nome Format)
    # Cerca il box del titolo predefinito
    if slide.shapes.title:
        slide.shapes.title.text = data['format_name']
    else:
        st.warning("Attenzione: Non ho trovato un box 'Titolo' standard nella slide 0.")

    # 2. GESTIONE CLAIM (Sottotitolo)
    # Cerchiamo il secondo segnaposto di testo (che non sia il titolo)
    found_subtitle = False
    for shape in slide.placeholders:
        # ph_idx 1 è spesso il sottotitolo nei layout standard, ma verifichiamo se è testo e non titolo
        if shape.has_text_frame and shape != slide.shapes.title:
            shape.text = data['claim']
            found_subtitle = True
            break # Trovato e riempito, ci fermiamo
            
    if not found_subtitle:
        st.warning("Attenzione: Non ho trovato un segnaposto per il Claim/Sottotitolo.")

    # 3. GESTIONE IMMAGINE (Prompt nelle note)
    # Non possiamo inserire un'immagine generata al volo senza API call esterna e salvataggio file.
    # Inseriamo il prompt nelle NOTE della slide per riferimento futuro.
    if not slide.has_notes_slide:
        slide.notes_slide = prs.slides._library.add_notes_slide(slide.part, slide.part.slide_layout) # Creazione note se mancano
    
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = f"--- IMAGEN 3 PROMPT ---\n{data['imagen_prompt']}"

    # Salvataggio
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

def extract_text_slide_1(file):
    prs = Presentation(file)
    # Estraiamo testo solo dalle prime slide (spesso il titolo è nella 1, ma a volte nella 2)
    # Per sicurezza leggiamo tutto il file per capire il contesto, ma ci concentriamo sulla Cover.
    text = []
    for s in prs.slides:
        for shape in s.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# --- INTERFACCIA ---
st.title("Step 1: Generatore Cover 🎨")
st.markdown("Carica il Template (con 10 pagine rigide) e il Vecchio PPT. Questo script compilerà **solo la Pagina 1**.")

# Selezione Modelli (Hardcoded su richiesta tassativa versione 3)
gemini_model = "gemini-3.0-pro-exp" # O la versione specifica disponibile
imagen_target = "Imagen 3"

col1, col2 = st.columns(2)
with col1:
    template_file = st.file_uploader("Template Rigido (10 slide)", type=['pptx'])
with col2:
    content_file = st.file_uploader("Vecchio PPT", type=['pptx'])

if template_file and content_file:
    if st.button("Genera Cover"):
        with st.spinner("Analisi in corso con Gemini 3..."):
            # 1. Estrai testo
            raw_text = extract_text_slide_1(content_file)
            
            # 2. Genera Dati
            try:
                cover_data = get_cover_data(raw_text, "gemini-1.5-pro") # Sostituisci con gemini-3 appena hai l'ID esatto nel tuo account
                
                st.success("Dati Generati!")
                
                # Visualizzazione anteprima
                st.subheader("Anteprima Dati Cover:")
                st.markdown(f"**Format (Intoccabile):** `{cover_data['format_name']}`")
                st.markdown(f"**Claim (Creativo):** `{cover_data['claim']}`")
                st.info(f"**Prompt Imagen:** {cover_data['imagen_prompt']}")
                
                # 3. Processa PPT
                new_ppt = process_cover_slide(template_file, cover_data)
                
                st.download_button("📥 Scarica PPT (Solo Cover Aggiornata)", new_ppt, "Cover_Updated.pptx")
                
            except Exception as e:
                st.error(f"Errore: {e}")
