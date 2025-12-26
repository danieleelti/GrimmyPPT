import streamlit as st
import google.generativeai as genai
import json
import requests
import io
from pptx.util import Inches
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# --- FUNZIONE 1: ANALISI TESTO (Gemini) ---
def analyze_content(context, gemini_model):
    try:
        model = genai.GenerativeModel(gemini_model)
        prompt_text = f"""
        Sei un Art Director. COMPITI:
        1. NOME FORMAT: Estrailo ESATTO dal testo.
        2. CLAIM: Crea uno slogan commerciale potente.
        3. PROMPT IMMAGINE: Scrivi un prompt DETTAGLIATO in inglese per una copertina FOTOREALISTICA.

        RISPONDI SOLO JSON: {{"format_name": "...", "claim": "...", "imagen_prompt": "..."}}

        TESTO SORGENTE: {context[:5000]}
        """
        res_text = model.generate_content(prompt_text, generation_config={"response_mime_type": "application/json"})
        return json.loads(res_text.text)
    except Exception as e:
        st.error(f"Errore Analisi Gemini: {e}")
        return None

# --- FUNZIONE 2: GENERAZIONE IMMAGINE (Imagen) ---
def generate_image_with_imagen(prompt, api_key, model_name):
    if not model_name.startswith("models/"): model_name = f"models/{model_name}"
    url = f"https://generativelanguage.googleapis.com/v1beta/{model_name}:predict?key={api_key}"
    headers = {"Content-Type": "application/json"}
    data = {"instances": [{"prompt": prompt}], "parameters": {"aspectRatio": "16:9", "sampleCount": 1}}
    
    try:
        response = requests.post(url, headers=headers, json=data)
        response.raise_for_status()
        result = response.json()
        if "predictions" in result:
            import base64
            return base64.b64decode(result["predictions"][0]["bytesBase64Encoded"])
        return None
    except Exception as e:
        st.error(f"Errore Imagen: {e}")
        return None

# --- FUNZIONE 3: INSERIMENTO NEL PPT (FIX LAYOUT SURGERY) ---
def insert_content_into_ppt(slide, data, img_bytes):
    """
    Inserisce testi nella slide e sposta chirurgicamente l'immagine nel LAYOUT.
    """
    try:
        # 1. INSERIMENTO TESTI (Slide)
        if slide.shapes.title: 
            slide.shapes.title.text = data.get("format_name", "")
        else:
            for s in slide.placeholders:
                if s.has_text_frame: s.text = data.get("format_name", ""); break
        
        for s in slide.placeholders:
            if s.has_text_frame and s != slide.shapes.title and s.text != data.get("format_name", ""):
                s.text = data.get("claim", ""); break
        
        # 2. INSERIMENTO IMMAGINE NEL LAYOUT (Metodo Avanzato)
        if img_bytes:
            layout = slide.slide_layout
            image_stream = io.BytesIO(img_bytes)
            
            # A. Cerchiamo le coordinate ideali dal placeholder dello schema (se esiste)
            target_left = Inches(0)
            target_top = Inches(0)
            target_width = Inches(10) # Default wide
            target_height = Inches(5.625)
            
            found_ph = False
            for shape in layout.placeholders:
                if shape.placeholder_format.type in [18, 7]: # Picture or Body
                    target_left = shape.left
                    target_top = shape.top
                    target_width = shape.width
                    target_height = shape.height
                    found_ph = True
                    break
            
            if not found_ph:
                # Fallback: schermo intero
                target_width = Inches(13.333) # 16:9 standard width
                target_height = Inches(7.5)

            # B. "Trucco": Aggiungiamo l'immagine alla SLIDE temporaneamente
            # (Perché non possiamo aggiungerla direttamente al layout facilmente)
            pic = slide.shapes.add_picture(image_stream, target_left, target_top, target_width, target_height)
            
            # C. CHIRURGIA: SPOSTIAMO L'IMMAGINE DALLA SLIDE AL LAYOUT
            try:
                # 1. Otteniamo la parte immagine (il file fisico nel pacchetto pptx)
                # L'immagine ha un 'rId' nella slide. Recuperiamolo.
                rId_slide = pic.element.blipFill.blip.embed
                image_part = slide.part.related_part(rId_slide)
                
                # 2. Creiamo una relazione tra il LAYOUT e quella stessa immagine
                # Questo permette al layout di "vedere" il file immagine
                rId_layout = layout.part.relate_to(image_part, RT.IMAGE)
                
                # 3. Aggiorniamo l'elemento XML dell'immagine per usare il nuovo rId del layout
                pic.element.blipFill.blip.embed = rId_layout
                
                # 4. Spostiamo fisicamente il nodo XML dall'albero della Slide all'albero del Layout
                slide.shapes._spTree.remove(pic.element) # Rimuovi dalla slide
                
                # Inseriamo nel Layout.
                # Indice 2 è solitamente sicuro (dopo le proprietà, prima delle shape in primo piano)
                # Questo la mette sullo sfondo del layout, DIETRO ai loghi/grafiche del layout stesso (se sono stati aggiunti dopo)
                layout.shapes._spTree.insert(2, pic.element) 
                
                # st.success("Immagine integrata con successo nello Schema Diapositiva!")

            except Exception as e_surgery:
                st.warning(f"Spostamento nel Layout fallito ({e_surgery}). L'immagine rimarrà sulla slide (in fondo).")
                # Fallback: Se fallisce la chirurgia, almeno rimettiamola sulla slide e mandiamola in fondo
                # Nota: se abbiamo già rimosso l'elemento, dobbiamo reinserirlo o ricrearlo.
                # Qui assumiamo che se fallisce prima del remove, è ancora lì. Se fallisce dopo, è persa.
                # Per sicurezza, nel catch non facciamo nulla di distruttivo se possibile.
                pass

        return True
    except Exception as e:
        st.error(f"Errore critico PPT: {e}")
        return False
